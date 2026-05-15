"""font-mcp 구성 도식 PPTX 생성기.

Usage:
    python scripts/make_architecture_pptx.py [output.pptx]

기본 출력: ./font-mcp-architecture.pptx

슬라이드:
    1) 표지
    2) 한눈에 보기 — 5-file 아키텍처
    3) MCP 도구 카탈로그
    4) 작업서 파이프라인 (Confluence → 폰트 산출물)
    5) XLSX 6-시트 스키마
    6) 학습 루프 (CaseMemory + ChromaDB)
    7) 데이터 / 디렉터리 레이아웃
"""
from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt


# ── 색 팔레트 ────────────────────────────────────────────────────
NAVY = RGBColor(0x2F, 0x5B, 0x7C)
BLUE = RGBColor(0x5B, 0x8F, 0xB9)
LIGHT = RGBColor(0xE8, 0xEE, 0xF4)
ACCENT = RGBColor(0xD9, 0x73, 0x4E)   # 주황 (작업서 강조)
GREEN = RGBColor(0x4F, 0x8A, 0x55)
GRAY_TEXT = RGBColor(0x37, 0x41, 0x51)
GRAY_BG = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x10, 0x18, 0x22)

FONT_KO = "맑은 고딕"


# ── 헬퍼 ─────────────────────────────────────────────────────────

def _set_text(
    tf,
    text: str,
    *,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = GRAY_TEXT,
    align: int = PP_ALIGN.LEFT,
    anchor: int = MSO_ANCHOR.TOP,
):
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    para = tf.paragraphs[0]
    para.alignment = align
    para.text = text
    for run in para.runs:
        run.font.name = FONT_KO
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def _add_lines(
    tf, lines: list[tuple[str, int, bool, RGBColor]],
    align: int = PP_ALIGN.LEFT,
):
    """lines: (text, size, bold, color)."""
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.text = text
        for run in p.runs:
            run.font.name = FONT_KO
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


def _box(
    slide, left, top, width, height, *,
    fill: RGBColor = WHITE,
    line: RGBColor = NAVY,
    line_width: float = 1.0,
    shape=MSO_SHAPE.ROUNDED_RECTANGLE,
):
    s = slide.shapes.add_shape(shape, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(line_width)
    s.shadow.inherit = False
    return s


def _label(
    slide, left, top, width, height, text, *,
    fill: RGBColor = WHITE,
    line: RGBColor = NAVY,
    text_color: RGBColor = GRAY_TEXT,
    size: int = 12,
    bold: bool = False,
    shape=MSO_SHAPE.ROUNDED_RECTANGLE,
    align=PP_ALIGN.CENTER,
):
    s = _box(slide, left, top, width, height, fill=fill, line=line, shape=shape)
    _set_text(
        s.text_frame, text,
        size=size, bold=bold, color=text_color,
        align=align, anchor=MSO_ANCHOR.MIDDLE,
    )
    return s


def _textbox(slide, left, top, width, height, text, **kwargs):
    tb = slide.shapes.add_textbox(left, top, width, height)
    _set_text(tb.text_frame, text, **kwargs)
    return tb


def _arrow(slide, x1, y1, x2, y2, *, color: RGBColor = NAVY, weight: float = 1.5):
    """Right arrow connector."""
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(weight)
    c.line.end_arrow_type = 2  # MSO_LINE_DASH_STYLE / arrow head
    # python-pptx 의 EndArrow 는 직접 XML 패치
    from pptx.oxml.ns import qn
    ln = c.line._get_or_add_ln()
    tail = ln.find(qn("a:tailEnd"))
    if tail is None:
        from lxml import etree
        tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")
    return c


def _slide_title(slide, text: str):
    tb = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.25), Inches(12.33), Inches(0.6),
    )
    _set_text(
        tb.text_frame, text,
        size=24, bold=True, color=NAVY, align=PP_ALIGN.LEFT,
    )
    # 밑줄 효과: 가는 라인
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(0.95), Inches(12.33), Emu(6350),  # ~0.007"
    )
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()


# ── 슬라이드 빌더 ────────────────────────────────────────────────

def _slide_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # 배경 강조 띠
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(0), Inches(13.33), Inches(2.6),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    _textbox(
        slide, Inches(0.7), Inches(0.7), Inches(12), Inches(1.2),
        "font-mcp 구성 도식",
        size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT,
    )
    _textbox(
        slide, Inches(0.7), Inches(1.7), Inches(12), Inches(0.6),
        "fontTools · ChromaDB · FastMCP 기반 폰트 제작 보조 서버",
        size=18, color=LIGHT, align=PP_ALIGN.LEFT,
    )

    # 핵심 요약 카드 3개
    cards = [
        ("입력", "Confluence [제작의뢰서]\n→ XLSX 작업서 (6-시트)",
         BLUE),
        ("처리", "FastMCP 23-tools\nFontOps + FontBuilder",
         NAVY),
        ("출력", "OTF / TTF / WOFF2 /\nVF · 자동 검증 + 학습",
         GREEN),
    ]
    x = Inches(0.7)
    for title, body, color in cards:
        card = _box(
            slide, x, Inches(3.2), Inches(4), Inches(2.4),
            fill=WHITE, line=color, line_width=2.0,
        )
        _textbox(
            slide, x + Inches(0.2), Inches(3.35), Inches(3.6), Inches(0.5),
            title, size=14, bold=True, color=color,
        )
        _textbox(
            slide, x + Inches(0.2), Inches(3.9), Inches(3.6), Inches(1.6),
            body, size=14, color=GRAY_TEXT,
        )
        x += Inches(4.2)

    _textbox(
        slide, Inches(0.7), Inches(6.8), Inches(12), Inches(0.4),
        "Sandoll Inc. · font-mcp v0.2 · 2026",
        size=11, color=GRAY_TEXT, align=PP_ALIGN.LEFT,
    )


def _slide_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_title(slide, "한눈에 보기 — 5-file 아키텍처")

    # Agent
    agent = _label(
        slide, Inches(0.5), Inches(2.6), Inches(1.8), Inches(1.0),
        "Claude\n(Agent)",
        fill=GRAY_BG, line=GRAY_TEXT, size=13, bold=True,
    )
    # MCP Server (server.py)
    server = _label(
        slide, Inches(3.0), Inches(2.6), Inches(2.6), Inches(1.0),
        "server.py\nFastMCP — 23 tools",
        fill=NAVY, line=NAVY, text_color=WHITE, size=13, bold=True,
    )
    # FontOps
    fontops = _label(
        slide, Inches(6.4), Inches(1.4), Inches(2.8), Inches(0.9),
        "font_ops.py\nFontOps — fontTools 래퍼",
        fill=WHITE, line=BLUE, size=12, bold=True,
    )
    # FontBuilder
    fb = _label(
        slide, Inches(6.4), Inches(2.6), Inches(2.8), Inches(0.9),
        "font_builder.py\nFontBuilder — 패밀리 빌드",
        fill=WHITE, line=ACCENT, size=12, bold=True,
    )
    # BuildSpec
    bs = _label(
        slide, Inches(6.4), Inches(3.8), Inches(2.8), Inches(0.9),
        "build_spec.py\nXLSX → BuildSpec 파서",
        fill=WHITE, line=ACCENT, size=12, bold=True,
    )
    # CaseMemory
    mem = _label(
        slide, Inches(6.4), Inches(5.0), Inches(2.8), Inches(0.9),
        "memory.py\nCaseMemory (RAG)",
        fill=WHITE, line=GREEN, size=12, bold=True,
    )
    # Chroma
    chroma = _label(
        slide, Inches(10.0), Inches(5.0), Inches(2.7), Inches(0.9),
        "ChromaDB\nall-MiniLM-L6-v2",
        fill=GRAY_BG, line=GREEN, size=12, bold=True,
    )
    # fontTools
    ft = _label(
        slide, Inches(10.0), Inches(1.4), Inches(2.7), Inches(0.9),
        "fontTools / fontbakery",
        fill=GRAY_BG, line=BLUE, size=12, bold=True,
    )

    # 화살표
    _arrow(slide, agent.left + agent.width, agent.top + Inches(0.5),
           server.left, server.top + Inches(0.5),
           color=GRAY_TEXT)
    for tgt in (fontops, fb, bs, mem):
        _arrow(slide,
               server.left + server.width, server.top + Inches(0.5),
               tgt.left, tgt.top + Inches(0.45),
               color=NAVY)
    _arrow(slide, fontops.left + fontops.width, fontops.top + Inches(0.45),
           ft.left, ft.top + Inches(0.45), color=BLUE)
    _arrow(slide, mem.left + mem.width, mem.top + Inches(0.45),
           chroma.left, chroma.top + Inches(0.45), color=GREEN)
    # 빌더가 ops + memory 둘 다 호출
    _arrow(slide,
           fb.left + Inches(1.4), fb.top,
           fontops.left + Inches(1.4), fontops.top + fontops.height,
           color=ACCENT)
    _arrow(slide,
           fb.left + Inches(1.4), fb.top + fb.height,
           mem.left + Inches(1.4), mem.top,
           color=ACCENT)

    # 캡션
    _textbox(
        slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.6),
        "서버는 항상 try/except 로 감싸 {\"error\": ...} 형태로 반환 — 에이전트가 다음 액션을 결정.",
        size=12, color=GRAY_TEXT,
    )


def _slide_tools(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_title(slide, "MCP 도구 카탈로그 — 23 tools")

    groups = [
        ("진단 (read-only)", BLUE, [
            "font_info", "list_tables", "dump_table_ttx",
            "get_name_records", "get_vertical_metrics",
            "list_features", "diagnose",
        ]),
        ("수정 (write — output_path 필수)", NAVY, [
            "apply_ttx_patch", "set_name_record", "set_vertical_metrics",
        ]),
        ("변환", BLUE, [
            "subset_font", "merge_fonts", "convert_format", "instance_variable",
        ]),
        ("시각 / 비교 / 검증", GREEN, [
            "render_sample", "diff_fonts", "validate_font",
        ]),
        ("학습 메모리 (RAG)", GREEN, [
            "find_similar_cases", "record_case",
            "update_case_outcome", "validate_and_record",
        ]),
        ("작업서 빌드 (XLSX → 패밀리)", ACCENT, [
            "parse_build_sheet", "build_font_family",
        ]),
    ]

    # 2 x 3 그리드
    cols, rows = 2, 3
    col_w = Inches(6.1)
    row_h = Inches(1.85)
    base_l = Inches(0.6)
    base_t = Inches(1.25)
    for i, (title, color, tools) in enumerate(groups):
        c, r = i % cols, i // cols
        left = base_l + col_w * c + Inches(0.05) * c
        top = base_t + row_h * r + Inches(0.1) * r

        # 헤더
        head = _box(
            slide, left, top, col_w, Inches(0.4),
            fill=color, line=color,
        )
        _set_text(
            head.text_frame, title,
            size=12, bold=True, color=WHITE,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
        )
        # 본문
        body = _box(
            slide, left, top + Inches(0.4), col_w, row_h - Inches(0.45),
            fill=WHITE, line=color,
        )
        # 칩처럼 한 줄에 콤마 구분
        text = "  ·  ".join(tools)
        _set_text(
            body.text_frame, text,
            size=12, color=GRAY_TEXT, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.MIDDLE,
        )


def _slide_pipeline(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_title(slide, "작업서 파이프라인 — Confluence → 폰트 산출물")

    # 6 단계 가로 흐름
    steps = [
        ("Confluence\n[제작의뢰서]", GRAY_BG, GRAY_TEXT),
        ("scripts/\nmake_build_template.py", LIGHT, NAVY),
        ("build_template.xlsx\n+ Google Sheets 편집", WHITE, ACCENT),
        ("parse_build_sheet\n→ BuildSpec", WHITE, NAVY),
        ("build_font_family\n(+ base_otf/ttf/VF)", WHITE, ACCENT),
        ("OTF / TTF / WOFF2\nWOFF2_subset / VF", GREEN, WHITE),
    ]
    n = len(steps)
    box_w = Inches(1.9)
    box_h = Inches(1.1)
    gap = Inches(0.2)
    total_w = box_w * n + gap * (n - 1)
    start_l = (Inches(13.33) - total_w) / 2
    top = Inches(2.5)

    last_right = None
    for i, (text, fill, line) in enumerate(steps):
        left = start_l + (box_w + gap) * i
        text_color = WHITE if fill == GREEN else GRAY_TEXT
        b = _label(
            slide, left, top, box_w, box_h, text,
            fill=fill, line=line, text_color=text_color,
            size=11, bold=True,
        )
        if last_right is not None:
            _arrow(
                slide, last_right, top + Inches(0.55),
                left, top + Inches(0.55), color=NAVY, weight=1.5,
            )
        last_right = left + box_w

    # 위 / 아래 보조 정보
    _textbox(
        slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
        "📥 입력: 디자이너 의뢰서  →  ⚙️ 빌드: 메트릭 + nameID 일괄 적용  →  📤 산출: 검증 + 학습",
        size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER,
    )

    _textbox(
        slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.4),
        "각 산출물은 fontbakery 결과(PASS/WARN/FAIL)와 함께 CaseMemory 에 자동 저장",
        size=12, color=GRAY_TEXT, align=PP_ALIGN.CENTER,
    )

    # 패치 항목 (per weight)
    patches_box = _box(
        slide, Inches(0.7), Inches(4.8), Inches(12), Inches(2.0),
        fill=LIGHT, line=NAVY,
    )
    _add_lines(
        patches_box.text_frame,
        [
            ("📋 build_font_family 가 weight 별로 적용하는 패치", 13, True, NAVY),
            ("", 6, False, NAVY),
            ("• head:  fontRevision · macStyle(BOLD/ITALIC)", 12, False, GRAY_TEXT),
            ("• OS/2:  usWeightClass · fsType · sTypo* · usWin* · strikeout* · fsSelection(REGULAR/BOLD/ITALIC + USE_TYPO_METRICS)", 12, False, GRAY_TEXT),
            ("• hhea / post:  ascent · descent · lineGap · underline*", 12, False, GRAY_TEXT),
            ("• name:  1·2·4·16·17 자동 (en + ko) · 3(uniqueID) · 5(version) · 6(PSname+suffix) · 0/7/8/9/11-14 (names 시트)", 12, False, GRAY_TEXT),
        ],
    )


def _slide_xlsx(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_title(slide, "XLSX 작업서 — 6-시트 스키마")

    sheets = [
        ("meta", "프로젝트 · 벤더 · 저작권",
         "project_name · version · vendor · designer_ko/en · vendor_url · license · copyright · trademark · name_record_scope(korean|latin)"),
        ("metrics", "head · hhea · OS/2 · post",
         "unitsPerEm · hhea ascent/descent/lineGap · sTypo* · usWin* · strikeout · fsType · underlinePosition/Thickness · use_typo_metrics"),
        ("weights", "웨이트별 행 (다중 행)",
         "weight_class · style_name · is_bold · is_italic · korean_family · latin_family · psname · fullname_suffix · base_font_key"),
        ("outputs", "산출 포맷 토글",
         "OTF · TTF · WOFF · WOFF2 · WOFF_subset · WOFF2_subset · VF  (each: enabled + psname_suffix)"),
        ("names", "추가 nameID 행",
         "nameID(0,7,8,9,11-14) · platform(win|mac) · lang(en|ko|...) · value"),
        ("subset", "서브셋 규칙",
         "mode(preset|text|unicodes) · preset=common_kr · layout_features(* 또는 콤마구분)"),
    ]

    # 2x3
    base_l = Inches(0.6)
    base_t = Inches(1.3)
    box_w = Inches(6.1)
    box_h = Inches(1.7)
    gap_x = Inches(0.15)
    gap_y = Inches(0.2)
    for i, (name, subtitle, fields) in enumerate(sheets):
        c, r = i % 2, i // 2
        left = base_l + (box_w + gap_x) * c
        top = base_t + (box_h + gap_y) * r

        # 좌측 시트명 띠
        side = _box(
            slide, left, top, Inches(1.4), box_h,
            fill=NAVY, line=NAVY,
        )
        _set_text(
            side.text_frame, name,
            size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # 우측 본문
        body = _box(
            slide, left + Inches(1.4), top,
            box_w - Inches(1.4), box_h,
            fill=WHITE, line=NAVY,
        )
        _add_lines(
            body.text_frame,
            [
                (subtitle, 12, True, NAVY),
                ("", 4, False, NAVY),
                (fields, 11, False, GRAY_TEXT),
            ],
        )


def _slide_learning(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_title(slide, "학습 루프 — CaseMemory + ChromaDB")

    # 중앙: ChromaDB
    chroma = _label(
        slide, Inches(5.1), Inches(3.0), Inches(3.1), Inches(1.2),
        "ChromaDB (collection: cases)\nall-MiniLM-L6-v2 · 로컬, no API key",
        fill=NAVY, line=NAVY, text_color=WHITE, size=12, bold=True,
    )

    # 왼쪽: 쓰기 흐름
    write_box = _label(
        slide, Inches(0.6), Inches(1.5), Inches(4.0), Inches(0.5),
        "✍️  쓰기 (Write)",
        fill=ACCENT, line=ACCENT, text_color=WHITE, size=13, bold=True,
        align=PP_ALIGN.LEFT,
    )
    write_steps = [
        ("set_vertical_metrics / set_name_record / apply_ttx_patch",
         WHITE, ACCENT),
        ("validate_font (fontbakery)\n→ PASS / WARN / FAIL", WHITE, ACCENT),
        ("record_case  /  validate_and_record\nsymptom · diagnosis · patch_ttx · verdict",
         LIGHT, ACCENT),
        ("build_font_family — weight 별 자동 기록",
         LIGHT, ACCENT),
    ]
    top = Inches(2.15)
    for text, fill, line in write_steps:
        _label(
            slide, Inches(0.6), top, Inches(4.0), Inches(0.6), text,
            fill=fill, line=line, size=10, bold=False,
        )
        top += Inches(0.7)

    # 오른쪽: 읽기 흐름
    read_box = _label(
        slide, Inches(8.7), Inches(1.5), Inches(4.0), Inches(0.5),
        "🔍  읽기 (Read)",
        fill=GREEN, line=GREEN, text_color=WHITE, size=13, bold=True,
        align=PP_ALIGN.LEFT,
    )
    read_steps = [
        ("diagnose(font) → issues[*].message", WHITE, GREEN),
        ("find_similar_cases(message, k=5)\n→ 임베딩 유사도 검색", WHITE, GREEN),
        ("score = success - fail\n검증된 패치가 상단", LIGHT, GREEN),
        ("apply_ttx_patch(patch_ttx)\nupdate_case_outcome(success=...)",
         LIGHT, GREEN),
    ]
    top = Inches(2.15)
    for text, fill, line in read_steps:
        _label(
            slide, Inches(8.7), top, Inches(4.0), Inches(0.6), text,
            fill=fill, line=line, size=10,
        )
        top += Inches(0.7)

    # 화살표: 쓰기 → ChromaDB ← 읽기
    _arrow(slide, Inches(4.6), Inches(3.6), Inches(5.1), Inches(3.6),
           color=ACCENT, weight=2.0)
    _arrow(slide, Inches(8.7), Inches(3.6), Inches(8.2), Inches(3.6),
           color=GREEN, weight=2.0)

    # 하단: 케이스 스키마
    schema = _box(
        slide, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.4),
        fill=GRAY_BG, line=GRAY_TEXT,
    )
    _add_lines(
        schema.text_frame,
        [
            ("📦 Case 스키마", 13, True, NAVY),
            ("", 4, False, NAVY),
            ("id · ts · symptom (document/임베딩) · diagnosis · patch_summary · table_tag · font_path",
             11, False, GRAY_TEXT),
            ("validation_after(PASS/WARN/FAIL) · patch_ttx(재적용용, 20KB cap) · success_count · fail_count · score",
             11, False, GRAY_TEXT),
        ],
    )


def _slide_data_layout(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_title(slide, "데이터 / 디렉터리 레이아웃")

    # 트리 박스
    tree = _box(
        slide, Inches(0.6), Inches(1.3), Inches(6.0), Inches(5.3),
        fill=GRAY_BG, line=NAVY,
    )
    _add_lines(
        tree.text_frame,
        [
            ("$FONT_MCP_DATA_DIR  (기본: <repo>/data)", 13, True, NAVY),
            ("  ├─ chroma/                   ChromaDB 영속", 11, False, GRAY_TEXT),
            ("  │    └─ chroma.sqlite3        cases 컬렉션", 11, False, GRAY_TEXT),
            ("  ├─ cases.json                레거시 케이스 (옵션)", 11, False, GRAY_TEXT),
            ("  ├─ wiki_cache/               Confluence 캐시", 11, False, GRAY_TEXT),
            ("  │    ├─ <page_id>.html        body.storage XHTML", 11, False, GRAY_TEXT),
            ("  │    └─ <page_id>.json/.title 메타", 11, False, GRAY_TEXT),
            ("  └─ font/                     샘플 베이스 폰트", 11, False, GRAY_TEXT),
            ("", 8, False, NAVY),
            ("<repo>/", 13, True, NAVY),
            ("  ├─ server.py                 MCP 엔트리", 11, False, GRAY_TEXT),
            ("  ├─ font_ops.py · memory.py", 11, False, GRAY_TEXT),
            ("  ├─ build_spec.py · font_builder.py  ⬅ 신규", 11, True, ACCENT),
            ("  ├─ build_template.xlsx       생성된 작업서 샘플 ⬅ 신규",
             11, True, ACCENT),
            ("  └─ scripts/", 11, False, GRAY_TEXT),
            ("        ├─ fetch_confluence.py  (페이지 → wiki_cache/)",
             11, False, GRAY_TEXT),
            ("        ├─ import_cases.py", 11, False, GRAY_TEXT),
            ("        └─ make_build_template.py  ⬅ 신규",
             11, True, ACCENT),
        ],
    )

    # 우측: 환경 변수 / Setup
    setup = _box(
        slide, Inches(6.9), Inches(1.3), Inches(5.8), Inches(2.5),
        fill=WHITE, line=NAVY,
    )
    _add_lines(
        setup.text_frame,
        [
            ("⚙️  환경 변수 (.env)", 14, True, NAVY),
            ("", 4, False, NAVY),
            ("FONT_MCP_DATA_DIR=~/.font-mcp/", 12, False, GRAY_TEXT),
            ("ATLASSIAN_EMAIL=…", 12, False, GRAY_TEXT),
            ("ATLASSIAN_TOKEN=…", 12, False, GRAY_TEXT),
            ("", 6, False, NAVY),
            ("• 셸 env 가 .env 보다 우선 (load_dotenv override=False)",
             10, False, GRAY_TEXT),
            ("• Atlassian 토큰: id.atlassian.com/manage-profile/security/api-tokens",
             10, False, GRAY_TEXT),
        ],
    )

    # 우측 아래: 실행
    run_box = _box(
        slide, Inches(6.9), Inches(4.0), Inches(5.8), Inches(2.6),
        fill=WHITE, line=ACCENT,
    )
    _add_lines(
        run_box.text_frame,
        [
            ("🚀  실행 흐름", 14, True, ACCENT),
            ("", 4, False, NAVY),
            ("1.  python scripts/make_build_template.py", 11, False, GRAY_TEXT),
            ("2.  Google Drive 업로드 → Sheets 편집 → .xlsx 다운로드",
             11, False, GRAY_TEXT),
            ("3.  python scripts/fetch_confluence.py <page_id>  (필요시)",
             11, False, GRAY_TEXT),
            ("4.  python server.py   (MCP stdio 시작)", 11, False, GRAY_TEXT),
            ("5.  Claude: parse_build_sheet → build_font_family(...)",
             11, False, GRAY_TEXT),
            ("6.  output_dir/{OTF,TTF,WOFF2,...} 산출 + 자동 검증/학습",
             11, False, GRAY_TEXT),
        ],
    )


# ── 엔트리 ───────────────────────────────────────────────────────

def build(out_path: Path) -> None:
    prs = Presentation()
    # 16:9 와이드
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    _slide_cover(prs)
    _slide_architecture(prs)
    _slide_tools(prs)
    _slide_pipeline(prs)
    _slide_xlsx(prs)
    _slide_learning(prs)
    _slide_data_layout(prs)

    prs.save(out_path)
    print(f"wrote {out_path}  (slides: {len(prs.slides)})")


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument(
        "output", nargs="?",
        default=str(
            Path(__file__).resolve().parent.parent
            / "font-mcp-architecture.pptx"
        ),
    )
    args = ap.parse_args()
    build(Path(args.output))


if __name__ == "__main__":
    main()
